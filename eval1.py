import streamlit as st
import tempfile, zipfile, os, json, httpx, re
from pdfminer.high_level import extract_text
import docx2txt
from pptx import Presentation
from langchain_openai import ChatOpenAI

# Disable SSL verification for internal TCS GenAI endpoint
client = httpx.Client(verify=False)

# ------------------------- Streamlit UI Setup -------------------------
st.set_page_config(page_title="AI Hackathon Evaluation Engine", page_icon="🤖", layout="wide")
st.title("🤖 Automated Hackathon Evaluation Engine (TCS GenAI Lab)")

st.markdown("""
This AI system automatically evaluates hackathon submissions by analyzing:
- 🧠 Code (`.py`)
- 📄 Documentation (`.pdf` / `.docx`)
- 🎞️ Presentations (`.pptx`)
- 🗒️ Additional files (`.txt`, `.md`)

The evaluation is based on predefined scoring rubrics such as **Relevance**, **Innovation**, **Feasibility**, **GenAI Usage**, and **Presentation Quality**.
""")

# ------------------------- Inputs -------------------------
api_key = st.text_input("🔑 Enter your TCS GenAI API Key", type="password")
rubric_text = st.text_area("📘 Paste or write your Scoring Rubric", height=180, placeholder="""
Example:
Relevance to the Problem and Tech Landscape - 30
Innovation and Creativity - 15
Feasibility and Scalability - 15
Usage of GenAI and AI/ML Techniques - 20
Teamwork and Presentation Quality - 20
""")
upload_file = st.file_uploader("📦 Upload a ZIP containing all submission files", type=["zip"])

# ------------------------- File Extraction -------------------------
def extract_submission_text(upload_file):
    import fitz  # PyMuPDF for better PDF extraction
    submission_text = ""

    with tempfile.TemporaryDirectory() as tmp:
        with zipfile.ZipFile(upload_file, "r") as zf:
            zf.extractall(tmp)

        for root, _, files in os.walk(tmp):
            for f in files:
                path = os.path.join(root, f)
                ext = f.lower().split(".")[-1]

                try:
                    if ext == "py":
                        submission_text += f"\n\n[FILE: {f} - CODE]\n" + open(path, encoding="utf-8").read()

                    elif ext == "pdf":
                        submission_text += f"\n\n[FILE: {f} - DOCUMENT]\n"
                        try:
                            text = extract_text(path)
                            if len(text.strip()) < 20:
                                raise ValueError("Low text extraction")
                        except Exception:
                            # Fallback to PyMuPDF for scanned PDFs
                            doc = fitz.open(path)
                            text = "".join(page.get_text("text") for page in doc)
                        submission_text += text

                    elif ext == "docx":
                        submission_text += f"\n\n[FILE: {f} - DOCUMENT]\n" + docx2txt.process(path)

                    elif ext == "pptx":
                        submission_text += f"\n\n[FILE: {f} - SLIDES]\n"
                        prs = Presentation(path)
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    submission_text += shape.text + "\n"

                    elif ext in ("txt", "md"):
                        submission_text += f"\n\n[FILE: {f}]\n" + open(path, encoding="utf-8").read()

                    elif ext in ("mp4", "mov", "mkv"):
                        submission_text += f"\n\n[FILE: {f} - VIDEO]\n(Video file detected — textual content not extracted.)"

                    else:
                        submission_text += f"\n\n[FILE: {f} - SKIPPED]\nUnsupported file type."

                except Exception as e:
                    submission_text += f"\n\n[FILE: {f} - ERROR]\nCould not read this file ({e})"

    return submission_text.strip()

# ------------------------- Evaluation Logic -------------------------
def evaluate_submission(api_key, rubric_text, submission_text):
    llm = ChatOpenAI(
        base_url="https://genailab.tcs.in",
        model="azure/genailab-maas-gpt-4o",
        api_key=api_key,
        http_client=client
    )

    prompt = f"""
You are an impartial AI evaluator for hackathon submissions.
also respond/dont 
Evaluate the following submission according to this rubric:
{rubric_text}

Submission:
{submission_text[:20000]}

Return a JSON in this exact format:
{{
  "criteria": [
    {{"name": "Relevance to the Problem and Tech Landscape", "score": <0-30>, "reason": "<why>"}},
    {{"name": "Innovation and Creativity", "score": <0-15>, "reason": "<why>"}},
    {{"name": "Feasibility and Scalability", "score": <0-15>, "reason": "<why>"}},
    {{"name": "Usage of GenAI and AI/ML Techniques", "score": <0-20>, "reason": "<why>"}},
    {{"name": "Teamwork and Presentation Quality", "score": <0-20>, "reason": "<why>"}}
  ],
  "total_score": <sum_of_scores_out_of_100>,
  "summary": "<1-2 line justification>"
}}
    """

    response = llm.invoke(prompt)
    return response

# ------------------------- Streamlit Workflow -------------------------
if api_key and rubric_text and upload_file:
    with st.spinner("📦 Extracting and analyzing submission files..."):
        submission_text = extract_submission_text(upload_file)

    st.success("✅ Files successfully extracted and processed.")
    with st.expander("🔍 Preview Extracted Text", expanded=False):
        st.text_area("Extracted Submission Content", submission_text[:4000], height=300)

    if st.button("⚖️ Run Evaluation"):
        with st.spinner("🤖 Evaluating submission with TCS GPT-4o..."):
            result = evaluate_submission(api_key, rubric_text, submission_text)

        # Extract JSON only from LLM response
        raw = getattr(result, "content", None) or str(result)
        json_match = re.search(r"\{[\s\S]*\}", raw)
        parsed = None
        if json_match:
            try:
                parsed = json.loads(json_match.group(0))
            except json.JSONDecodeError:
                parsed = None

        if not parsed:
            st.warning("⚠️ Model output was not pure JSON. Showing raw output.")
            st.write(raw)

        if parsed:
            # Normalize and clean numeric formatting
            total_score = parsed.get("total_score", "N/A")
            try:
                if isinstance(total_score, (int, float)) and total_score <= 10:
                    total_score = total_score * 10
            except:
                pass

            st.subheader("📊 Evaluation Summary")
            st.metric("Final Score", f"{round(total_score, 1)}/100")

            st.markdown("### Detailed Criteria Scores")
            for c in parsed.get("criteria", []):
                score = c["score"]
                try:
                    if isinstance(score, (int, float)) and score <= 10:
                        score = round(score * 3, 1)
                except:
                    pass
                st.write(f"**{c['name']}**: {score} — {c['reason']}")

            st.markdown("### 🧠 Overall Summary")
            st.write(parsed.get("summary", ""))

            # Download JSON report
            json_report = json.dumps(parsed, indent=2)
            st.download_button("📥 Download Evaluation Report (JSON)", json_report, file_name="evaluation_report.json")

else:
    st.info("👆 Please enter your API key, paste rubric, and upload ZIP to start evaluation.")
