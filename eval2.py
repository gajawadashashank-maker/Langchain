import streamlit as st
import tempfile, zipfile, os, json, httpx, re, time
from pdfminer.high_level import extract_text
import docx2txt
from pptx import Presentation
from langchain_openai import ChatOpenAI

# Disable SSL verification for internal TCS GenAI endpoint
client = httpx.Client(verify=False)

# ------------------------- Streamlit UI Setup -------------------------
st.set_page_config(page_title="Multi-Team Hackathon Evaluation Engine", page_icon="🤖", layout="wide")
st.title("🤖 Multi-Team Automated Hackathon Evaluation Engine (TCS GenAI Lab)")

st.markdown("""
This AI system automatically evaluates **multiple hackathon submissions** by analyzing:
- 🧠 Code (`.py`)
- 📄 Documentation (`.pdf` / `.docx`)
- 🎞️ Presentations (`.pptx`)
- 🗒️ Text notes (`.txt`, `.md`)

It applies a predefined **rubric** (Relevance, Innovation, Feasibility, GenAI Usage, Presentation)  
and generates a **scoreboard** for all submissions.
""")

# ------------------------- Inputs -------------------------
api_key = st.text_input("🔑 Enter your TCS GenAI API Key", type="password")
rubric_text = st.text_area("📘 Paste or write your Scoring Rubric", height=160, placeholder="""
Relevance to the Problem and Tech Landscape - 30
Innovation and Creativity - 15
Feasibility and Scalability - 15
Usage of GenAI and AI/ML Techniques - 20
Teamwork and Presentation Quality - 20
""")

uploads = st.file_uploader("📦 Upload one or more ZIP files (each team’s submission)", type=["zip"], accept_multiple_files=True)

# ------------------------- Helper: Extract Text -------------------------
def extract_submission_text(upload_file):
    import fitz  # PyMuPDF for PDFs
    submission_text = ""

    with tempfile.TemporaryDirectory() as tmp:
        with zipfile.ZipFile(upload_file, "r") as zf:
            zf.extractall(tmp)

        for root, _, files in os.walk(tmp):
            for f in files:
                path = os.path.join(root, f)
                ext = f.lower().split(".")[-1]
                try:
                    if ext == "py":
                        submission_text += f"\n\n[FILE: {f} - CODE]\n" + open(path, encoding="utf-8").read()
                    elif ext == "pdf":
                        submission_text += f"\n\n[FILE: {f} - DOCUMENT]\n"
                        try:
                            text = extract_text(path)
                            if len(text.strip()) < 20:
                                raise ValueError("Low text extraction")
                        except Exception:
                            doc = fitz.open(path)
                            text = "".join(page.get_text("text") for page in doc)
                        submission_text += text
                    elif ext == "docx":
                        submission_text += f"\n\n[FILE: {f} - DOCUMENT]\n" + docx2txt.process(path)
                    elif ext == "pptx":
                        submission_text += f"\n\n[FILE: {f} - SLIDES]\n"
                        prs = Presentation(path)
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    submission_text += shape.text + "\n"
                    elif ext in ("txt", "md"):
                        submission_text += f"\n\n[FILE: {f}]\n" + open(path, encoding="utf-8").read()
                    elif ext in ("mp4", "mov", "mkv"):
                        submission_text += f"\n\n[FILE: {f} - VIDEO]\n(Video detected – textual content skipped.)"
                except Exception as e:
                    submission_text += f"\n\n[FILE: {f} - ERROR]\nCould not read this file ({e})"
    return submission_text.strip()

# ------------------------- LLM Evaluation -------------------------
def evaluate_submission(api_key, rubric_text, submission_text):
    llm = ChatOpenAI(
        base_url="https://genailab.tcs.in",
        model="azure/genailab-maas-gpt-4o",
        api_key=api_key,
        http_client=client
    )

    prompt = f"""
You are an impartial AI evaluator for hackathon submissions.

Evaluate the following submission using the given rubric:
{rubric_text}

Submission:
{submission_text[:20000]}

Return a JSON in this exact format:
{{
  "criteria": [
    {{"name": "Relevance to the Problem and Tech Landscape", "score": <0-30>, "reason": "<why>"}},
    {{"name": "Innovation and Creativity", "score": <0-15>, "reason": "<why>"}},
    {{"name": "Feasibility and Scalability", "score": <0-15>, "reason": "<why>"}},
    {{"name": "Usage of GenAI and AI/ML Techniques", "score": <0-20>, "reason": "<why>"}},
    {{"name": "Teamwork and Presentation Quality", "score": <0-20>, "reason": "<why>"}}
  ],
  "total_score": <sum_out_of_100>,
  "summary": "<short justification>"
}}
    """

    result = llm.invoke(prompt)
    raw = getattr(result, "content", None) or str(result)
    json_match = re.search(r"\{[\s\S]*\}", raw)
    if json_match:
        try:
            return json.loads(json_match.group(0))
        except:
            pass
    return {"error": "Invalid JSON", "raw": raw}

# ------------------------- Multi-Team Evaluation Flow -------------------------
if api_key and rubric_text and uploads:
    st.success(f"✅ {len(uploads)} submission(s) uploaded. Ready for evaluation.")
    if st.button("⚖️ Evaluate All Submissions"):
        all_results = []
        progress = st.progress(0)
        for i, upload_file in enumerate(uploads, start=1):
            st.write(f"🧩 Evaluating: `{upload_file.name}` ...")
            try:
                submission_text = extract_submission_text(upload_file)
                result = evaluate_submission(api_key, rubric_text, submission_text)
                if "total_score" in result:
                    total = result["total_score"]
                else:
                    total = 0
                all_results.append({
                    "Team Name": upload_file.name.replace(".zip", ""),
                    "Score": total,
                    "Summary": result.get("summary", "No summary"),
                    "Details": result
                })
            except Exception as e:
                all_results.append({
                    "Team Name": upload_file.name,
                    "Score": 0,
                    "Summary": f"Error: {e}",
                    "Details": {}
                })
            progress.progress(i / len(uploads))
            time.sleep(1)

        # ----------------- Leaderboard -----------------
        st.subheader("🏆 Evaluation Leaderboard")
        leaderboard = sorted(all_results, key=lambda x: x["Score"], reverse=True)
        for rank, entry in enumerate(leaderboard, start=1):
            st.markdown(f"### 🥇 Rank {rank}: {entry['Team Name']}")
            st.metric("Final Score", f"{entry['Score']}/100")
            st.write(f"🧠 Summary: {entry['Summary']}")
            with st.expander("📊 View Detailed Criteria"):
                st.json(entry["Details"])
            st.markdown("---")

        # ----------------- Downloadable Report -----------------
        st.download_button(
            "📥 Download All Evaluations (JSON)",
            json.dumps(leaderboard, indent=2),
            file_name="all_evaluations.json"
        )

else:
    st.info("👆 Please enter API key, paste rubric, and upload ZIPs to start multi-team evaluation.")
